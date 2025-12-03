"""Generate a Title Lift project presentation deck."""

from __future__ import annotations

from datetime import UTC, datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


OUTPUT_PATH = Path("docs/presentations/title_lift_project_brief.pptx")
OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)


def _add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Title Lift Pipeline"
    subtitle = slide.placeholders[1]
    subtitle.text = "Separating intrinsic quality from headline lift\nPatrick Ray · {}".format(
        datetime.now(UTC).strftime("%Y-%m-%d")
    )


def _add_bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for idx, bullet in enumerate(bullets):
        paragraph = tf.add_paragraph() if idx else tf.paragraphs[0]
        paragraph.text = bullet
        paragraph.font.size = Pt(20)
        paragraph.level = 0


def _add_split_slide(prs: Presentation, title: str, left_items: list[str], right_items: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[3])  # Title and two content columns
    slide.shapes.title.text = title
    left_box = slide.shapes.placeholders[1].text_frame
    left_box.clear()
    for idx, item in enumerate(left_items):
        paragraph = left_box.add_paragraph() if idx else left_box.paragraphs[0]
        paragraph.text = item
        paragraph.font.size = Pt(18)

    right_box = slide.shapes.placeholders[2].text_frame
    right_box.clear()
    for idx, item in enumerate(right_items):
        paragraph = right_box.add_paragraph() if idx else right_box.paragraphs[0]
        paragraph.text = item
        paragraph.font.size = Pt(18)


def build_deck() -> Presentation:
    prs = Presentation()
    _add_title_slide(prs)

    _add_bullets_slide(
        prs,
        "Agenda",
        [
            "Motivation & proposal alignment",
            "Data collection & preprocessing",
            "Stage A intrinsic modeling",
            "Stage B title-lift modeling",
            "Diagnostics & results",
            "Implications and next steps",
        ],
    )

    _add_bullets_slide(
        prs,
        "Research Motivation",
        [
            "Extend Weissburg et al. (ICWSM 2022) with refreshed Reddit + Hacker News corpus.",
            "Question: How much incremental lift do titles provide once early exposure is controlled?",
            "Success bar: maintain stable intrinsic RMSE while achieving ≥60% pairwise ranking accuracy.",
        ],
    )

    _add_bullets_slide(
        prs,
        "Pipeline Overview",
        [
            "Collectors (`bin/collect_reddit.py`, `bin/collect_hn.py`) gather score snapshots at 5–60 minutes.",
            "`bin/make_features.py` normalizes platforms, hashes authors, and engineers sentiment/clickbait features.",
            "Notebook + Stage scripts consume `data/features.parquet` for modeling and diagnostics.",
        ],
    )

    _add_split_slide(
        prs,
        "Data Inventory",
        [
            "13.4k posts across technology, business, world news communities.",
            "Score snapshots: 5m / 15m / 30m / 60m for early velocity.",
            "Feature blocks: sentiment, readability, clickbait heuristics, entity counts.",
        ],
        [
            "Context features: hour/day buckets, author & subreddit frequency, content-type flags.",
            "Human-readable export: `docs/title_lift_human_readable.csv` (47 curated columns).",
            "Artifacts logged under `outputs/title_lift/` for reproducibility.",
        ],
    )

    _add_bullets_slide(
        prs,
        "Stage A – Intrinsic Model",
        [
            "Log-linear OLS on `log1p(score_60m)` using exposure (`log1p(score_5m)`) + context covariates.",
            "Calibration checks: residual histograms, QQ plots, hourly/subreddit bias tables.",
            "Delivers ~58% variance explained on November holdout; residuals feed Stage B.",
        ],
    )

    _add_bullets_slide(
        prs,
        "Stage B – Title Lift",
        [
            "OLS residual regression on title features (length, sentiment, clickbait heuristics).",
            "Mean-centered, Elastic Net, and LightGBM/SHAP variants validate signal stability.",
            "Outputs include coefficient tables, feature importances, residual scatter plots.",
        ],
    )

    _add_split_slide(
        prs,
        "Diagnostics Stack",
        [
            "Pairwise ranking evaluation (within subreddit/hour) to measure ≥60% success target.",
            "Temporal quantile splits & blocked CV highlight robustness across time windows.",
            "Bootstrap + learning curves quantify uncertainty and data sufficiency.",
        ],
        [
            "Residual dashboards: hour-of-day drift, platform boxplots, residual vs fitted scatter.",
            "Rollup tables compare Stage A vs Stage B variants (OLS, Elastic Net, LightGBM).",
            "Figures exported to `docs/figures/` for manuscript inclusion.",
        ],
    )

    _add_bullets_slide(
        prs,
        "Key Results",
        [
            "Stage B reduces RMSE modestly while improving ranking accuracy toward the 0.60 target.",
            "Largest title-lift gains in technology/business subreddits; moderation-heavy domains remain flat.",
            "Human-readable dataset + final report package the findings for faculty review.",
        ],
    )

    _add_bullets_slide(
        prs,
        "Implications",
        [
            "Headline experimentation is justified where evening cohorts and upbeat sentiment prevail.",
            "Exposure proxies still leak in certain off-hours; need richer context features for production.",
            "Automation hooks ready for Stage B deployment once monitoring is established.",
        ],
    )

    _add_bullets_slide(
        prs,
        "Next Steps",
        [
            "Automate regular data refresh + QA gates for sentiment and clickbait dictionaries.",
            "Revisit contextual embeddings once compute budget is restored for Stage B enhancements.",
            "Design online A/B tests to validate ≥60% pairwise lift in live recommendation flows.",
        ],
    )

    _add_bullets_slide(
        prs,
        "Artifacts & Submission",
        [
            "Notebook: `Title_Lift_Pipeline.ipynb` (reproducible modeling + diagnostics).",
            "Report: `docs/title_lift_final_report.docx` (publishable manuscript draft).",
            "Datasets: `data/features.parquet` (analysis) + `docs/title_lift_human_readable.csv` (submission).",
        ],
    )

    return prs


if __name__ == "__main__":
    presentation = build_deck()
    presentation.save(OUTPUT_PATH)
    print(f"Presentation written to {OUTPUT_PATH}")
